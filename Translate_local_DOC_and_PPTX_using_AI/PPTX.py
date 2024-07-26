from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re

def process_shapes(shapes, process_function, progress_callback=None, total_shapes=None, shape_counter=None):
    for shape in shapes:
        print(f"Processing shape: {shape.name}")  # Added for progress tracking
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                if paragraph.text:
                    match = re.match(r"(\d+\.\s*)(.*)", paragraph.text)
                    if match:
                        prefix, text = match.groups()
                        print(f"Processing paragraph: {text[:20]}...")
                        translated_text = process_function(text)
                        translated_text = prefix + translated_text
                    else:
                        print(f"Processing paragraph: {paragraph.text[:40]}...")
                        translated_text = process_function(paragraph.text)

                    for i, run in enumerate(paragraph.runs):
                        if i == 0:
                            run.text = translated_text
                        else:
                            run.text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        if paragraph.text:
                            match = re.match(r"(\d+\.\s*)(.*)", paragraph.text)
                            if match:
                                prefix, text = match.groups()
                                print(f"Processing table cell paragraph: {text[:20]}...")
                                translated_text = process_function(text)
                                translated_text = prefix + translated_text
                            else:
                                print(f"Processing table cell paragraph: {paragraph.text[:20]}...")
                                translated_text = process_function(paragraph.text)

                            for i, run in enumerate(paragraph.runs):
                                if i == 0:
                                    run.text = translated_text
                                else:
                                    run.text = ""
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            print(f"Processing group: {shape.name}")
            process_shapes(shape.shapes, process_function, progress_callback, total_shapes, shape_counter)

        if progress_callback and total_shapes and shape_counter:
            shape_counter[0] += 1
            progress = shape_counter[0] / total_shapes * 100
            progress_callback(progress)

def process_slides(file_path, process_function, progress_callback=None):
    prs = Presentation(file_path)
    total_slides = len(prs.slides)
    total_shapes = sum(len(slide.shapes) for slide in prs.slides)
    shape_counter = [0]

    for index, slide in enumerate(prs.slides):
        process_shapes(slide.shapes, process_function, progress_callback, total_shapes, shape_counter)
        if progress_callback:
            progress = (index + 1) / total_slides * 100
            progress_callback(progress)

    return prs
